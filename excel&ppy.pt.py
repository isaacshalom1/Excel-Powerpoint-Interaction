import openpyxl 
from pptx import Presentation 

#Enter Excel file and powerpoints to search through 
wb = openpyxl.load_workbook("")
boys_ppt = Presentation("")
girls_ppt = Presentation("")

ws = wb.active

count = 2
for col in ws.iter_cols(max_col =1, min_row =2, max_row =ws.max_row): #Column 1 from first row to last
	for cell in col:
		#print("hello " + str(count))
		tag = cell.value #grab cell value 
		stop = False

		for slide in boys_ppt.slides: #for the slides
			if( not stop):
				for shape in slide.shapes: #for slide shapes
					if( not stop):
						try:
							if shape.has_text_frame:
								if shape.text.find(tag) != -1: #if you found tag 
								 	 title = slide.shapes.title.text
								 	 ws["B" +str(count)] = title #grab title enter it into enxcel row 
								 	 count +=1	 	 		#increment row	
								 	 stop = True #stop the iteration 
						except:
							print("found exception at " + str(count))


 	 		
		for slide2 in girls_ppt.slides:
			if(not stop):
				for shape2 in slide2.shapes:
					if(not stop):
				 		if shape2.has_text_frame:
				 	 		if shape2.text.find(tag) != -1:
				 	 			title = slide2.shapes.title.text
				 	 			ws["B" +str(count)] = title
				 	 			count+=1
				 	 			stop = True

		if stop == False: #if generic 
			ws["B" +str(count)] = ""
			count+=1
			
#enter excel file 
wb.save("")
			
